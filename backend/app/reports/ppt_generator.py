from pptx import Presentation
from datetime import datetime
from pathlib import Path

from app.models.stock_models import StockSummary


TEMPLATE_PATH = Path(__file__).parent / "templates" / "stock_report_template.pptx"
OUTPUT_DIR = Path(__file__).parent / "generated"


def generate_stock_report(stock: StockSummary) -> Path:

    OUTPUT_DIR.mkdir(exist_ok=True)

    presentation = Presentation(TEMPLATE_PATH)

    replacements = {
        "{{SYMBOL}}": stock.symbol,
        "{{DATE}}": datetime.now().strftime("%Y-%m-%d"),
        "{{COMPANY_NAME}}": stock.company.name,
        "{{EXCHANGE}}": stock.company.exchange,
        "{{INDUSTRY}}": stock.company.industry,
        "{{PRICE}}": str(stock.price),
        "{{MARKET_CAP}}": str(stock.company.market_cap)
    }

    for slide in presentation.slides:
        for shape in slide.shapes:

            # Handle normal text frames
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for key, value in replacements.items():
                            if key in run.text:
                                run.text = run.text.replace(key, value)

            # Handle tables while preserving formatting
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                for key, value in replacements.items():
                                    if key in run.text:
                                        run.text = run.text.replace(key, value)


    output_file = OUTPUT_DIR / f"{stock.symbol}_report.pptx"

    presentation.save(output_file)

    return output_file
